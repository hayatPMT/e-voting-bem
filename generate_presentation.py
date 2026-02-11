from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_evoting_presentation():
    """Create comprehensive E-Voting BEM system presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    BLUE_DARK = RGBColor(21, 101, 192)
    BLUE_LIGHT = RGBColor(66, 165, 245)
    GREEN = RGBColor(67, 160, 71)
    ORANGE = RGBColor(251, 140, 0)
    RED = RGBColor(229, 57, 53)
    GRAY = RGBColor(97, 97, 97)
    WHITE = RGBColor(255, 255, 255)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SISTEM E-VOTING BEM"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Dokumentasi Alur Proses & Arsitektur Sistem"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = BLUE_LIGHT
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Daftar Isi
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "DAFTAR ISI"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    topics = [
        "1. Gambaran Umum Sistem",
        "2. Arsitektur Database",
        "3. Peran Pengguna (User Roles)",
        "4. Alur Proses Mahasiswa",
        "5. Alur Proses Admin",
        "6. Alur Voting Lengkap",
        "7. Fitur-Fitur Sistem",
        "8. Keamanan & Validasi",
        "9. Teknologi yang Digunakan"
    ]
    
    for topic in topics:
        p = tf.add_paragraph()
        p.text = topic
        p.font.size = Pt(20)
        p.space_before = Pt(10)
        p.level = 0

    # Slide 3: Gambaran Umum Sistem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "GAMBARAN UMUM SISTEM"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    overview = [
        ("Tujuan", "Sistem pemilihan BEM secara digital yang aman dan transparan"),
        ("Platform", "Web-based application menggunakan Laravel 12"),
        ("Pengguna", "Admin dan Mahasiswa"),
        ("Fitur Utama", "Verifikasi mahasiswa, voting, real-time chart, manajemen data"),
        ("Keamanan", "Autentikasi, role-based access control, one-time voting")
    ]
    
    for label, desc in overview:
        p = tf.add_paragraph()
        p.text = f"{label}: "
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        
        run = p.runs[0]
        run2 = p.add_run()
        run2.text = desc
        run2.font.bold = False
        run2.font.color.rgb = GRAY
        
        p.space_before = Pt(15)

    # Slide 4: Arsitektur Database
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ARSITEKTUR DATABASE"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Tabel Utama:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    tables = [
        "1. USERS - Data autentikasi & role pengguna",
        "2. ADMIN_PROFILES - Profil lengkap admin",
        "3. MAHASISWA_PROFILES - Profil lengkap mahasiswa (NIM, jurusan, status voting)",
        "4. KANDIDATS - Data kandidat & pasangan calon",
        "5. VOTES - Rekaman suara yang masuk",
        "6. SETTINGS - Konfigurasi periode voting"
    ]
    
    for table in tables:
        p = tf.add_paragraph()
        p.text = table
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        p.level = 1
    
    # Add relationship info
    p = tf.add_paragraph()
    p.text = "\nRelasi:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_before = Pt(20)
    
    relations = [
        "• User (1:1) Admin Profile / Mahasiswa Profile",
        "• User (1:M) Votes",
        "• Kandidat (1:M) Votes"
    ]
    
    for rel in relations:
        p = tf.add_paragraph()
        p.text = rel
        p.font.size = Pt(16)
        p.space_before = Pt(5)

    # Slide 5: Peran Pengguna
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PERAN PENGGUNA (USER ROLES)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    # Create two columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_tf = left_box.text_frame
    
    p = left_tf.add_paragraph()
    p.text = "👨‍💼 ADMIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    admin_tasks = [
        "✓ Login dengan email & password",
        "✓ Manajemen kandidat (CRUD)",
        "✓ Manajemen mahasiswa (CRUD)",
        "✓ Import/export data mahasiswa",
        "✓ Manajemen admin lain",
        "✓ Lihat rekap hasil voting",
        "✓ Konfigurasi periode voting",
        "✓ Monitor status voting real-time"
    ]
    
    for task in admin_tasks:
        p = left_tf.add_paragraph()
        p.text = task
        p.font.size = Pt(14)
        p.space_before = Pt(8)
    
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    right_tf = right_box.text_frame
    
    p = right_tf.add_paragraph()
    p.text = "🎓 MAHASISWA"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    mhs_tasks = [
        "✓ Verifikasi dengan NIM & Password",
        "✓ Lihat daftar kandidat",
        "✓ Pilih kandidat (1x voting)",
        "✓ Konfirmasi pilihan",
        "✓ Lihat hasil voting real-time",
        "",
        "⚠️ Tidak bisa login ulang",
        "⚠️ Tidak bisa ubah pilihan"
    ]
    
    for task in mhs_tasks:
        p = right_tf.add_paragraph()
        p.text = task
        p.font.size = Pt(14)
        p.space_before = Pt(8)

    # Slide 6: Alur Proses Mahasiswa - Verifikasi
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ALUR MAHASISWA: VERIFIKASI"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    steps = [
        ("1", "Landing Page", "Mahasiswa mengakses halaman utama"),
        ("2", "Klik 'Mahasiswa'", "Diarahkan ke halaman verifikasi (/verifikasi)"),
        ("3", "Input NIM & Password", "Mahasiswa memasukkan kredensial"),
        ("4", "Sistem Validasi", "Cek NIM di database mahasiswa_profiles"),
        ("5", "Cek Status Voting", "Validasi apakah sudah pernah voting (has_voted)"),
        ("6", "Create Session", "Jika valid & belum voting → buat session"),
        ("7", "Redirect", "Arahkan ke halaman voting (/voting)")
    ]
    
    for num, step, desc in steps:
        p = tf.add_paragraph()
        p.text = f"{num}. {step}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_before = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.level = 1

    # Slide 7: Alur Mahasiswa - Voting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ALUR MAHASISWA: VOTING"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    steps = [
        ("1", "Halaman Voting", "Tampilkan semua kandidat dengan foto, visi, misi"),
        ("2", "Pilih Kandidat", "Mahasiswa klik tombol 'Vote' pada kandidat"),
        ("3", "Konfirmasi Modal", "Pop-up konfirmasi pilihan"),
        ("4", "Submit Vote", "Mahasiswa confirm → POST ke /vote/{id}"),
        ("5", "Sistem Proses", "Simpan ke tabel VOTES (user_id, kandidat_id)"),
        ("6", "Update Status", "Update mahasiswa_profiles: has_voted=true, voted_at=now()"),
        ("7", "Destroy Session", "Hapus session mahasiswa"),
        ("8", "Redirect Chart", "Arahkan ke halaman hasil (/chart)")
    ]
    
    for num, step, desc in steps:
        p = tf.add_paragraph()
        p.text = f"{num}. {step}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_before = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.level = 1

    # Slide 8: Alur Admin - Login & Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ALUR ADMIN: LOGIN & DASHBOARD"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    steps = [
        ("LOGIN", [
            "1. Akses halaman /login",
            "2. Input email & password admin",
            "3. Sistem validasi di tabel USERS (role='admin')",
            "4. Cek is_active = true",
            "5. Create authenticated session",
            "6. Update last_login timestamp",
            "7. Redirect ke /dashboard"
        ]),
        ("DASHBOARD", [
            "• Statistik total mahasiswa, kandidat, suara masuk",
            "• Grafik real-time hasil voting",
            "• Quick actions (tambah kandidat, mahasiswa)",
            "• Menu navigasi lengkap"
        ])
    ]
    
    for section, items in steps:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_before = Pt(15)
        
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = item
            p2.font.size = Pt(14)
            p2.space_before = Pt(5)
            p2.level = 1

    # Slide 9: Alur Admin - Manajemen Kandidat
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ALUR ADMIN: MANAJEMEN KANDIDAT"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    operations = [
        ("📋 LIHAT (INDEX)", [
            "/admin/kandidat → tampilkan semua kandidat",
            "Tabel dengan foto, nama, visi, misi, jumlah suara",
            "Tombol Edit & Delete per kandidat"
        ]),
        ("➕ TAMBAH (CREATE)", [
            "/admin/kandidat/create → form input",
            "Input: Nama, Visi, Misi, Upload Foto",
            "Validasi & simpan ke tabel KANDIDATS",
            "Upload foto ke storage/app/public/kandidat"
        ]),
        ("✏️ EDIT (UPDATE)", [
            "/admin/kandidat/{id}/edit → load data kandidat",
            "Update data kandidat yang dipilih",
            "Opsi ganti foto kandidat"
        ]),
        ("🗑️ HAPUS (DELETE)", [
            "Konfirmasi penghapusan",
            "Hapus foto dari storage",
            "Hapus record dari database"
        ])
    ]
    
    for op, details in operations:
        p = tf.add_paragraph()
        p.text = op
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_before = Pt(12)
        
        for detail in details:
            p2 = tf.add_paragraph()
            p2.text = f"  • {detail}"
            p2.font.size = Pt(12)
            p2.space_before = Pt(3)
            p2.level = 1

    # Slide 10: Alur Admin - Manajemen Mahasiswa
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ALUR ADMIN: MANAJEMEN MAHASISWA"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    operations = [
        ("📋 LIHAT & FILTER", [
            "Tabel semua mahasiswa dengan NIM, nama, prodi, angkatan",
            "Status voting (sudah/belum)",
            "Filter & search functionality"
        ]),
        ("➕ TAMBAH MANUAL", [
            "Form input data mahasiswa lengkap",
            "Auto-generate user account (email, password)",
            "Create mahasiswa_profile"
        ]),
        ("📥 IMPORT CSV", [
            "Download template CSV",
            "Upload file CSV dengan data mahasiswa bulk",
            "Validasi & insert ke database",
            "Report hasil import (success/failed)"
        ]),
        ("📤 EXPORT CSV", [
            "Export semua data mahasiswa ke CSV",
            "Include status voting & timestamp"
        ]),
        ("🔄 TOGGLE STATUS", [
            "Enable/disable mahasiswa dari voting",
            "Update field voting_enabled"
        ])
    ]
    
    for op, details in operations:
        p = tf.add_paragraph()
        p.text = op
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_before = Pt(10)
        
        for detail in details:
            p2 = tf.add_paragraph()
            p2.text = f"  • {detail}"
            p2.font.size = Pt(11)
            p2.space_before = Pt(2)
            p2.level = 1

    # Slide 11: Alur Voting Lengkap (Flowchart)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    title = slide.shapes.title
    title.text = "ALUR VOTING LENGKAP"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    # Create flowchart text
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf = flow_box.text_frame
    
    flow_steps = [
        "START → Landing Page",
        "  ↓",
        "Mahasiswa Klik 'Mahasiswa' → /verifikasi",
        "  ↓",
        "Input NIM & Password → Submit Form",
        "  ↓",
        "Validasi Kredensial & Status Voting",
        "  ↓                    ↓",
        "VALID               INVALID",
        "  ↓                    ↓",
        "Login Success       Error Message",
        "Session Created     (NIM salah / sudah voting)",
        "  ↓",
        "/voting → Tampil Daftar Kandidat",
        "  ↓",
        "Pilih Kandidat → Konfirmasi",
        "  ↓",
        "Submit Vote → Simpan ke DB",
        "  ↓",
        "Update has_voted=true, voted_at=now()",
        "  ↓",
        "Destroy Session",
        "  ↓",
        "Redirect ke /chart → Lihat Hasil",
        "  ↓",
        "END"
    ]
    
    for step in flow_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.space_before = Pt(2)
        if "→" in step or "START" in step or "END" in step:
            p.font.bold = True
            p.font.color.rgb = ORANGE

    # Slide 12: Fitur Keamanan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "KEAMANAN & VALIDASI"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RED
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    security_features = [
        ("🔒 AUTENTIKASI", [
            "Password di-hash menggunakan bcrypt",
            "Session-based authentication",
            "CSRF protection di semua form"
        ]),
        ("🛡️ AUTORISASI", [
            "Role-based access control (admin/mahasiswa)",
            "Middleware protection untuk route admin",
            "Mahasiswa hanya bisa voting jika sudah verifikasi"
        ]),
        ("✅ VALIDASI VOTING", [
            "Cek has_voted sebelum tampilkan halaman voting",
            "Prevent double voting (database constraint)",
            "Session dihapus setelah voting",
            "Mahasiswa tidak bisa login ulang ke sistem voting"
        ]),
        ("⏰ PERIODE VOTING", [
            "Admin set voting_start & voting_end",
            "Sistem auto-disable voting di luar periode",
            "Validasi timestamp sebelum terima vote"
        ]),
        ("📊 INTEGRITAS DATA", [
            "Foreign key constraints",
            "Unique constraints (NIM, email)",
            "Transaction untuk operasi critical"
        ])
    ]
    
    for feature, details in security_features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RED
        p.space_before = Pt(10)
        
        for detail in details:
            p2 = tf.add_paragraph()
            p2.text = f"  • {detail}"
            p2.font.size = Pt(12)
            p2.space_before = Pt(2)
            p2.level = 1

    # Slide 13: Fitur-Fitur Sistem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "FITUR-FITUR SISTEM"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_tf = left_box.text_frame
    
    p = left_tf.add_paragraph()
    p.text = "UNTUK MAHASISWA"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    mhs_features = [
        "✓ Verifikasi NIM cepat & aman",
        "✓ UI voting yang intuitif",
        "✓ Preview kandidat (foto, visi, misi)",
        "✓ Konfirmasi sebelum submit vote",
        "✓ Real-time chart hasil voting",
        "✓ Responsive design (mobile-friendly)",
        "✓ One-time voting (tidak bisa ganda)",
        "✓ Notifikasi sukses voting"
    ]
    
    for feature in mhs_features:
        p = left_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    right_tf = right_box.text_frame
    
    p = right_tf.add_paragraph()
    p.text = "UNTUK ADMIN"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    admin_features = [
        "✓ Dashboard analytics",
        "✓ CRUD kandidat lengkap",
        "✓ CRUD mahasiswa + import CSV",
        "✓ Export data mahasiswa",
        "✓ Manajemen admin accounts",
        "✓ Real-time vote monitoring",
        "✓ Rekap hasil voting",
        "✓ Konfigurasi periode voting",
        "✓ Toggle voting status mahasiswa",
        "✓ Activity logging"
    ]
    
    for feature in admin_features:
        p = right_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.space_before = Pt(6)

    # Slide 14: Public Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "FITUR PUBLIC (TANPA LOGIN)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "📊 HALAMAN CHART (/chart)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_before = Pt(10)
    
    features = [
        "• Accessible tanpa login",
        "• Real-time chart dengan Chart.js",
        "• Auto-refresh setiap beberapa detik",
        "• Menampilkan hasil voting sementara",
        "• Bar chart atau pie chart",
        "• Responsive layout",
        "• Bisa di-share ke publik"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "\n🎯 API ENDPOINT"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "GET /api/chart"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.space_before = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "Return: JSON dengan labels (nama kandidat) & values (jumlah suara)"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # Slide 15: Teknologi yang Digunakan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "TEKNOLOGI YANG DIGUNAKAN"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech_stack = [
        ("🔧 BACKEND", [
            "PHP 8.5",
            "Laravel Framework 12",
            "MySQL Database",
            "Eloquent ORM"
        ]),
        ("🎨 FRONTEND", [
            "HTML5, CSS3, JavaScript",
            "Vanilla CSS (responsive design)",
            "Chart.js (visualisasi data)",
            "Modern UI/UX design"
        ]),
        ("🔐 KEAMANAN", [
            "Laravel Authentication",
            "Bcrypt password hashing",
            "CSRF protection",
            "Session management"
        ]),
        ("🛠️ TOOLS & LIBRARY", [
            "Composer (dependency management)",
            "NPM (frontend packages)",
            "Laravel Pint (code formatting)",
            "Pest PHP (testing)"
        ])
    ]
    
    for category, items in tech_stack:
        p = tf.add_paragraph()
        p.text = category
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_before = Pt(12)
        
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = f"  • {item}"
            p2.font.size = Pt(14)
            p2.space_before = Pt(3)
            p2.level = 1

    # Slide 16: Struktur Route
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "STRUKTUR ROUTE"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    # Left column - Public & Mahasiswa
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_tf = left_box.text_frame
    
    p = left_tf.add_paragraph()
    p.text = "PUBLIC ROUTES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    public_routes = [
        "/ → Landing page",
        "/chart → Public chart",
        "/api/chart → Chart API"
    ]
    
    for route in public_routes:
        p = left_tf.add_paragraph()
        p.text = route
        p.font.size = Pt(11)
        p.space_before = Pt(4)
    
    p = left_tf.add_paragraph()
    p.text = "\nMAHASISWA ROUTES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_before = Pt(15)
    
    mhs_routes = [
        "/verifikasi → Verifikasi NIM",
        "/voting → Halaman voting",
        "/vote/{id} → Submit vote"
    ]
    
    for route in mhs_routes:
        p = left_tf.add_paragraph()
        p.text = route
        p.font.size = Pt(11)
        p.space_before = Pt(4)
    
    # Right column - Admin
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    right_tf = right_box.text_frame
    
    p = right_tf.add_paragraph()
    p.text = "ADMIN ROUTES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    admin_routes = [
        "/login → Admin login",
        "/logout → Logout",
        "/dashboard → Dashboard",
        "/rekap → Rekap hasil",
        "",
        "KANDIDAT:",
        "/admin/kandidat → List",
        "/admin/kandidat/create → Form",
        "/admin/kandidat/{id}/edit → Edit",
        "",
        "MAHASISWA:",
        "/admin/mahasiswa → List",
        "/admin/mahasiswa/create → Form",
        "/admin/mahasiswa/import → Import",
        "/admin/mahasiswa/export → Export",
        "",
        "SETTINGS:",
        "/admin/settings → Config",
        "/admin/admins → Admin mgmt"
    ]
    
    for route in admin_routes:
        p = right_tf.add_paragraph()
        if route == "":
            p.text = ""
        elif ":" in route:
            p.text = route
            p.font.bold = True
            p.font.size = Pt(12)
        else:
            p.text = route
            p.font.size = Pt(10)
        p.space_before = Pt(3)

    # Slide 17: Database Relationships Detail
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "RELASI DATABASE DETAIL"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    relationships = [
        ("USERS ↔ PROFILES (1:1)", [
            "• Satu user hanya punya 1 profil",
            "• Admin → admin_profiles (via user_id)",
            "• Mahasiswa → mahasiswa_profiles (via user_id)",
            "• On delete: CASCADE (hapus user = hapus profile)"
        ]),
        ("USERS ↔ VOTES (1:M)", [
            "• Satu user (mahasiswa) bisa punya banyak votes",
            "• TAPI di-constraint hanya 1 vote per user",
            "• Foreign key: votes.user_id → users.id"
        ]),
        ("KANDIDAT ↔ VOTES (1:M)", [
            "• Satu kandidat bisa terima banyak votes",
            "• Foreign key: votes.kandidat_id → kandidats.id",
            "• Count votes untuk dapat total suara per kandidat"
        ])
    ]
    
    for rel, details in relationships:
        p = tf.add_paragraph()
        p.text = rel
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_before = Pt(15)
        
        for detail in details:
            p2 = tf.add_paragraph()
            p2.text = detail
            p2.font.size = Pt(14)
            p2.space_before = Pt(5)
            p2.level = 1

    # Slide 18: Validasi & Error Handling
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "VALIDASI & ERROR HANDLING"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RED
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    validations = [
        ("📝 FORM REQUEST VALIDATION", [
            "Semua input di-validasi dengan Form Request classes",
            "Custom error messages dalam Bahasa Indonesia",
            "Client-side + server-side validation"
        ]),
        ("⚠️ ERROR SCENARIOS", [
            "NIM tidak ditemukan → redirect dengan error message",
            "Password salah → tampilkan error",
            "Sudah voting → tidak bisa akses /voting",
            "Voting di luar periode → disable voting",
            "Kandidat tidak valid → validation error"
        ]),
        ("🔄 TRANSACTION HANDLING", [
            "Critical operations wrapped dalam transaction",
            "Rollback jika ada error",
            "Ensure data consistency"
        ]),
        ("📋 LOGGING", [
            "Activity log untuk admin actions",
            "Error logging untuk debugging",
            "Access log untuk security audit"
        ])
    ]
    
    for section, items in validations:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RED
        p.space_before = Pt(12)
        
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = f"  • {item}"
            p2.font.size = Pt(13)
            p2.space_before = Pt(4)
            p2.level = 1

    # Slide 19: Best Practices
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "BEST PRACTICES IMPLEMENTASI"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    practices = [
        "✓ MVC Architecture (Model-View-Controller)",
        "✓ Repository pattern untuk data access",
        "✓ Service classes untuk business logic",
        "✓ Eloquent ORM untuk database queries",
        "✓ Middleware untuk authorization",
        "✓ Form Request untuk validation",
        "✓ Resource Controllers untuk CRUD",
        "✓ Named routes untuk URL generation",
        "✓ Blade templates untuk views",
        "✓ Asset compilation dengan Vite",
        "✓ Environment-based configuration",
        "✓ Database seeding untuk testing",
        "✓ Migration untuk version control database",
        "✓ Responsive design (mobile-first)",
        "✓ SEO-friendly structure",
        "✓ Code formatting dengan Laravel Pint"
    ]
    
    for practice in practices:
        p = tf.add_paragraph()
        p.text = practice
        p.font.size = Pt(15)
        p.space_before = Pt(8)

    # Slide 20: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "POTENSI PENGEMBANGAN"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    enhancements = [
        ("🔐 KEAMANAN TAMBAHAN", [
            "Two-factor authentication (2FA)",
            "Email verification untuk mahasiswa",
            "Biometric verification (fingerprint/face)",
            "Rate limiting untuk prevent brute force"
        ]),
        ("📊 ANALYTICS LANJUTAN", [
            "Voting statistics by prodi, angkatan",
            "Participation rate monitoring",
            "Demographic analysis",
            "Export report PDF/Excel"
        ]),
        ("🚀 FITUR TAMBAHAN", [
            "Live streaming hasil voting",
            "Email notification setelah voting",
            "QR code untuk quick access",
            "Multiple voting sessions/periode",
            "Voting untuk berbagai posisi (Ketua, Wakil, dll)",
            "Komentar/feedback untuk kandidat"
        ]),
        ("📱 MOBILE APP", [
            "Native mobile app (iOS/Android)",
            "Push notifications",
            "Offline capability"
        ])
    ]
    
    for category, items in enhancements:
        p = tf.add_paragraph()
        p.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.space_before = Pt(12)
        
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = f"  • {item}"
            p2.font.size = Pt(12)
            p2.space_before = Pt(3)
            p2.level = 1

    # Slide 21: Closing Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "TERIMA KASIH"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    
    p = subtitle_frame.paragraphs[0]
    p.text = "SISTEM E-VOTING BEM"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    p2 = subtitle_frame.add_paragraph()
    p2.text = "\nSistem Pemilihan Digital yang Aman dan Transparan"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('E-VOTING_BEM_PRESENTATION.pptx')
    print("✅ Presentasi berhasil dibuat: E-VOTING_BEM_PRESENTATION.pptx")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_evoting_presentation()
